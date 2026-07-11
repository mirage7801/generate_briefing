# -*- coding: utf-8 -*-
"""
AI + hardware co-design vision briefing generator.

This script reads a PDF or TXT paper, chunks it with section-aware overlap,
uses an LLM to create a NIL (Named terms, Important actions, Literals/Data)
intermediate representation, extracts a final briefing JSON, and renders a
three-slide PowerPoint briefing with python-pptx.
"""

from __future__ import annotations

import argparse
import getpass
import json
import os
import re
import sys
import textwrap
import warnings
from pathlib import Path
from typing import Any

from openai import OpenAI
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from pypdf import PdfReader


DEFAULT_MODEL = "gpt-4o-mini"
_LLM_CLIENT: OpenAI | None = None
_LLM_MODEL = DEFAULT_MODEL

COLORS = {
    "deep_blue": RGBColor(0x0B, 0x1F, 0x3A),
    "tech_blue": RGBColor(0x1F, 0x77, 0xB4),
    "light_blue": RGBColor(0xEA, 0xF4, 0xFF),
    "deep_gray": RGBColor(0x2B, 0x2B, 0x2B),
    "mid_gray": RGBColor(0x6B, 0x72, 0x80),
    "light_gray": RGBColor(0xF3, 0xF6, 0xFA),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
}


class BriefingError(RuntimeError):
    """Raised for user-facing briefing generation failures."""


def get_llm_config(args: argparse.Namespace) -> tuple[OpenAI, str]:
    """Build an OpenAI-compatible client and select a model without persisting secrets."""
    api_key = (
        getattr(args, "api_key", None)
        or os.getenv("API_KEY")
        or os.getenv("OPENAI_API_KEY")
        or os.getenv("DEEPSEEK_API_KEY")
    )
    if not api_key:
        api_key = getpass.getpass("No API key found. Please enter your API key: ")
    if not api_key:
        raise BriefingError("API key cannot be empty.")

    base_url = getattr(args, "base_url", None) or os.getenv("API_BASE_URL")
    model = getattr(args, "model", None) or os.getenv("API_MODEL") or DEFAULT_MODEL
    client = OpenAI(api_key=api_key, base_url=base_url) if base_url else OpenAI(api_key=api_key)
    return client, model


def load_document(file_path: str | Path) -> str:
    """Load text from a PDF or TXT document."""
    path = Path(file_path)
    if not path.exists():
        raise BriefingError(f"Input file does not exist: {path}")

    suffix = path.suffix.lower()
    if suffix == ".pdf":
        try:
            reader = PdfReader(str(path))
        except Exception as exc:
            raise BriefingError(f"Failed to open PDF file: {path}. {exc}") from exc

        page_texts: list[str] = []
        for idx, page in enumerate(reader.pages, start=1):
            try:
                extracted = page.extract_text() or ""
            except Exception as exc:
                warnings.warn(f"Page {idx} could not be extracted and was skipped: {exc}")
                continue
            if extracted.strip():
                page_texts.append(extracted)

        text = "\n\n".join(page_texts).strip()
        if not text:
            raise BriefingError(
                "No usable text could be extracted from the PDF. "
                "If this is a scanned PDF, run OCR first."
            )
        return text

    if suffix == ".txt":
        for encoding in ("utf-8", "gbk"):
            try:
                return path.read_text(encoding=encoding)
            except UnicodeDecodeError:
                continue
            except Exception as exc:
                raise BriefingError(f"Failed to read TXT file: {path}. {exc}") from exc
        raise BriefingError(f"Failed to decode TXT file with utf-8 or gbk: {path}")

    raise BriefingError(f"Unsupported input format: {suffix}. Please use PDF or TXT.")


def clean_text(text: str) -> str:
    """Normalize whitespace while preserving paragraph and section boundaries."""
    if not text:
        return ""

    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t\f\v]+", " ", text)
    text = re.sub(r" *\n *", "\n", text)

    raw_blocks = re.split(r"\n\s*\n+", text)
    cleaned_blocks: list[str] = []
    for block in raw_blocks:
        lines = [line.strip() for line in block.splitlines() if line.strip()]
        if not lines:
            continue

        if len(lines) == 1:
            cleaned_blocks.append(lines[0])
            continue

        joined_parts: list[str] = []
        for line in lines:
            if detect_section_title(line):
                if joined_parts:
                    cleaned_blocks.append(" ".join(joined_parts).strip())
                    joined_parts = []
                cleaned_blocks.append(line)
            else:
                joined_parts.append(line)
        if joined_parts:
            cleaned_blocks.append(" ".join(joined_parts).strip())

    return "\n\n".join(cleaned_blocks).strip()


def detect_section_title(paragraph: str) -> str:
    """Return a normalized section title if a paragraph looks like one."""
    text = paragraph.strip()
    if not text or len(text) > 120:
        return ""

    title_patterns = [
        r"^\d+(\.\d+)*\.?\s+[A-Z][A-Za-z0-9 ,:;()\-/&]+$",
        r"^[IVXLCDM]+\.\s+[A-Z][A-Za-z0-9 ,:;()\-/&]+$",
        r"^(Abstract|Introduction|Background|Related Work|Method|Methods|Approach|"
        r"Evaluation|Experiment|Experiments|Discussion|Conclusion|Conclusions|"
        r"References|Acknowledgements|Appendix)$",
    ]
    for pattern in title_patterns:
        if re.match(pattern, text, flags=re.IGNORECASE):
            return text

    words = text.split()
    if 1 <= len(words) <= 8 and text[:1].isupper():
        alpha_chars = re.sub(r"[^A-Za-z]", "", text)
        if alpha_chars and sum(ch.isupper() for ch in alpha_chars) / len(alpha_chars) > 0.65:
            return text

    return ""


def split_long_paragraph(paragraph: str, max_chars: int) -> list[str]:
    """Split a paragraph at sentence boundaries, falling back to safe character cuts."""
    paragraph = paragraph.strip()
    if len(paragraph) <= max_chars:
        return [paragraph]

    sentence_pattern = r"(?<=[。；;：:？！?!])\s+|(?<=[。；;：:？！?!])|(?<=\.)\s+"
    sentences = [s.strip() for s in re.split(sentence_pattern, paragraph) if s.strip()]
    if not sentences:
        sentences = [paragraph]

    parts: list[str] = []
    current = ""

    def safe_cut(text: str) -> list[str]:
        chunks: list[str] = []
        remaining = text.strip()
        while len(remaining) > max_chars:
            cut = remaining.rfind(" ", 0, max_chars)
            if cut < max_chars * 0.5:
                cut = max_chars
            chunks.append(remaining[:cut].strip())
            remaining = remaining[cut:].strip()
        if remaining:
            chunks.append(remaining)
        return chunks

    for sentence in sentences:
        if len(sentence) > max_chars:
            if current:
                parts.append(current.strip())
                current = ""
            parts.extend(safe_cut(sentence))
        elif not current:
            current = sentence
        elif len(current) + 1 + len(sentence) <= max_chars:
            current = f"{current} {sentence}"
        else:
            parts.append(current.strip())
            current = sentence

    if current:
        parts.append(current.strip())
    return parts


def chunk_text(text: str, max_chars: int = 3000, overlap_chars: int = 250) -> list[dict[str, Any]]:
    """Create section-aware, overlapped text chunks."""
    if max_chars <= 500:
        raise BriefingError("--max-chars should be greater than 500.")
    if overlap_chars < 0 or overlap_chars >= max_chars:
        raise BriefingError("--overlap must be non-negative and smaller than --max-chars.")

    paragraphs = [p.strip() for p in re.split(r"\n\s*\n+", text) if p.strip()]
    chunks: list[dict[str, Any]] = []
    current_parts: list[str] = []
    current_section_title = ""
    chunk_section_title = ""

    def build_text(parts: list[str], section_title: str) -> str:
        body = "\n\n".join(parts).strip()
        if section_title and not body.startswith(section_title):
            body = f"{section_title}\n\n{body}" if body else section_title
        return body

    def make_overlap(chunk_body: str, section_title: str) -> str:
        if overlap_chars <= 0 or not chunk_body:
            return ""
        tail = chunk_body[-overlap_chars:].strip()
        first_boundary = min(
            [pos for pos in [tail.find("\n\n"), tail.find(". "), tail.find("。")] if pos >= 0] or [-1]
        )
        if first_boundary >= 0 and first_boundary + 2 < len(tail):
            tail = tail[first_boundary + 2 :].strip()
        if section_title and section_title.lower() not in tail[:160].lower():
            tail = f"{section_title}\n\n{tail}"
        return tail.strip()

    def flush() -> str:
        nonlocal current_parts, chunk_section_title
        if not current_parts:
            return ""
        chunk_body = build_text(current_parts, chunk_section_title)
        chunks.append(
            {
                "chunk_id": len(chunks) + 1,
                "section_title": chunk_section_title or "",
                "text": chunk_body,
            }
        )
        overlap = make_overlap(chunk_body, chunk_section_title)
        current_parts = [overlap] if overlap else []
        return overlap

    for paragraph in paragraphs:
        title = detect_section_title(paragraph)
        if title:
            current_section_title = title

        pieces = split_long_paragraph(paragraph, max_chars)
        for piece in pieces:
            proposed_section = current_section_title or chunk_section_title
            if not current_parts:
                chunk_section_title = proposed_section

            proposed_parts = current_parts + [piece]
            proposed_text = build_text(proposed_parts, proposed_section)
            if current_parts and len(proposed_text) > max_chars:
                flush()
                chunk_section_title = proposed_section
                proposed_parts = current_parts + [piece]
                proposed_text = build_text(proposed_parts, proposed_section)

            current_parts = proposed_parts
            chunk_section_title = proposed_section

            if len(proposed_text) > max_chars * 1.15:
                flush()
                chunk_section_title = current_section_title

    if current_parts:
        body = build_text(current_parts, chunk_section_title)
        if body.strip():
            chunks.append(
                {
                    "chunk_id": len(chunks) + 1,
                    "section_title": chunk_section_title or "",
                    "text": body,
                }
            )

    return chunks


def build_nil_prompt(chunk: dict[str, Any]) -> str:
    """Build the per-chunk NIL extraction prompt."""
    return f"""
你是一名资深 AI+硬件协同设计技术分析师。

任务：从下面这个论文 chunk 中抽取 NIL 中间层信息。
NIL = Named terms + Important actions + Literals/Data。

要求：
1. 只输出严格 JSON，不输出 Markdown，不添加解释。
2. 不编造原文没有的信息。
3. 英文技术术语尽量保留英文。
4. 不只保留名词，也要保留动作关系和关键数据。
5. 如果某个字段没有内容，输出空数组。
6. chunk_id 和 section_title 必须与输入一致。

输出 JSON Schema：
{{
  "chunk_id": {chunk["chunk_id"]},
  "section_title": "{chunk.get("section_title", "")}",
  "named_terms": [
    {{
      "term": "...",
      "type": "技术术语/层级名称/论文目标/时间趋势/系统名称/其他",
      "context": "该术语在本 chunk 中的含义或作用"
    }}
  ],
  "important_actions": [
    {{
      "subject": "...",
      "action": "...",
      "object": "...",
      "meaning": "这组动作关系说明了什么"
    }}
  ],
  "literals_data": [
    {{
      "value": "...",
      "type": "倍数/时间范围/年份/指标/数量级/其他",
      "context": "该数据在论文中的含义"
    }}
  ],
  "candidate_insights": [
    "本 chunk 中可能与最终 PPT 相关的关键信息"
  ]
}}

论文 chunk：
{chunk["text"]}
""".strip()


def nil_preprocess(chunks: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Run NIL extraction on each chunk and keep chunk provenance."""
    nil_results: list[dict[str, Any]] = []
    schema_hint = "NIL JSON with chunk_id, section_title, named_terms, important_actions, literals_data, candidate_insights."

    for chunk in chunks:
        print(f"Running NIL preprocessing for chunk {chunk['chunk_id']}/{len(chunks)}...")
        raw = call_llm(build_nil_prompt(chunk))
        parsed = parse_or_repair_json(raw, schema_hint)
        parsed["chunk_id"] = chunk["chunk_id"]
        parsed["section_title"] = parsed.get("section_title") or chunk.get("section_title", "")
        nil_results.append(parsed)

    return nil_results


def build_merge_nil_prompt(nil_results: list[dict[str, Any]]) -> str:
    """Build the NIL merge prompt."""
    nil_json = json.dumps(nil_results, ensure_ascii=False, indent=2)
    return f"""
你是一名资深 AI+硬件协同设计技术分析师。

任务：合并多个 chunk 的 NIL 结果，去重并保留来源 chunk_id。

合并目标：
1. 合并重复术语、动作关系和数据。
2. 保留出现频率高、与最终简报强相关的信息。
3. 特别关注核心愿景、效率提升目标、Intelligence per Joule、Hardware / Algorithm / Application 三大层级、近期 2-5 年趋势、远期 6-10 年趋势。
4. 不添加 NIL 中不存在的新事实。

只输出严格 JSON，不输出 Markdown，不解释。

输出 JSON Schema：
{{
  "global_named_terms": [
    {{
      "term": "...",
      "type": "...",
      "merged_context": "...",
      "source_chunk_ids": [1, 2]
    }}
  ],
  "global_actions": [
    {{
      "subject": "...",
      "action": "...",
      "object": "...",
      "meaning": "...",
      "source_chunk_ids": [1, 2]
    }}
  ],
  "global_data": [
    {{
      "value": "...",
      "type": "...",
      "merged_context": "...",
      "source_chunk_ids": [1, 2]
    }}
  ],
  "high_confidence_insights": [
    {{
      "insight": "...",
      "source_chunk_ids": [1, 2]
    }}
  ]
}}

NIL 结果：
{nil_json}
""".strip()


def merge_nil_results(nil_results: list[dict[str, Any]]) -> dict[str, Any]:
    """Merge per-chunk NIL results into a global NIL representation."""
    raw = call_llm(build_merge_nil_prompt(nil_results))
    schema_hint = "Merged NIL JSON with global_named_terms, global_actions, global_data, high_confidence_insights."
    return parse_or_repair_json(raw, schema_hint)


def build_final_extraction_prompt(merged_nil: dict[str, Any], selected_context: str) -> str:
    """Build the final briefing extraction prompt."""
    merged_nil_json = json.dumps(merged_nil, ensure_ascii=False, indent=2)
    return f"""
你是一名资深 AI+硬件协同设计技术分析师和 PPT 简报架构师。

任务：基于 merged NIL 和 selected original context，生成最终 briefing JSON。

重要规则：
1. 输出中文，关键英文技术术语保留英文。
2. 只输出严格 JSON，不输出 Markdown，不解释。
3. 不编造论文没有的信息；论文没有明确说明时写“论文未明确说明”。
4. 最终内容必须优先基于 merged NIL；如果 selected context 与 NIL 冲突，以 selected context 为准。
5. layers 必须包含 Hardware、Algorithm、Application 三项。
6. timeline 必须包含近期和远期两项。
7. 不要把关键词当成答案，必须根据论文语义提取。

输出 JSON Schema：
{{
  "paper_title": "...",
  "core_vision": {{
    "headline": "...",
    "intelligence_per_joule": "...",
    "target_explanation": "..."
  }},
  "layers": [
    {{
      "name": "Hardware",
      "role": "...",
      "key_features": ["...", "...", "..."],
      "examples": ["...", "..."]
    }},
    {{
      "name": "Algorithm",
      "role": "...",
      "key_features": ["...", "...", "..."],
      "examples": ["...", "..."]
    }},
    {{
      "name": "Application",
      "role": "...",
      "key_features": ["...", "...", "..."],
      "examples": ["...", "..."]
    }}
  ],
  "timeline": [
    {{
      "stage": "近期",
      "range": "2-5 年",
      "trends": ["...", "...", "..."],
      "meaning": "..."
    }},
    {{
      "stage": "远期",
      "range": "6-10 年",
      "trends": ["...", "...", "..."],
      "meaning": "..."
    }}
  ],
  "one_sentence_summary": "..."
}}

merged NIL：
{merged_nil_json}

selected original context：
{selected_context}
""".strip()


def call_llm(prompt: str) -> str:
    """Call the configured OpenAI-compatible Chat Completions API."""
    if _LLM_CLIENT is None:
        raise BriefingError("LLM client is not configured.")

    try:
        response = _LLM_CLIENT.chat.completions.create(
            model=_LLM_MODEL,
            messages=[
                {
                    "role": "system",
                    "content": "You return only valid JSON when asked for JSON. Do not include Markdown.",
                },
                {"role": "user", "content": prompt},
            ],
            temperature=0.2,
        )
        content = response.choices[0].message.content or ""
    except Exception as exc:
        error_message = str(exc)
        configured_key = getattr(_LLM_CLIENT, "api_key", None)
        if configured_key:
            error_message = error_message.replace(configured_key, "[REDACTED]")
        raise BriefingError(f"LLM API call failed: {error_message}") from exc

    if not content.strip():
        raise BriefingError("LLM API returned empty output.")
    return content


def parse_json_strict(raw_output: str) -> Any:
    """Parse JSON, accepting accidental fenced code blocks but no extra prose."""
    text = raw_output.strip()
    text = re.sub(r"^```(?:json)?\s*", "", text, flags=re.IGNORECASE)
    text = re.sub(r"\s*```$", "", text)
    return json.loads(text)


def repair_json(raw_output: str, expected_schema_hint: str) -> Any:
    """Ask the LLM to repair invalid JSON without adding new information."""
    prompt = f"""
你只需要修复 JSON 格式。

规则：
1. 不要添加新信息。
2. 不要删除关键信息。
3. 不要输出 Markdown。
4. 不要解释。
5. 只输出合法 JSON。
6. 目标结构提示：{expected_schema_hint}

需要修复的原始输出：
{raw_output}
""".strip()
    repaired = call_llm(prompt)
    return parse_json_strict(repaired)


def parse_or_repair_json(raw_output: str, expected_schema_hint: str) -> Any:
    """Parse JSON; if parsing fails, ask the LLM to repair it."""
    try:
        return parse_json_strict(raw_output)
    except Exception as first_exc:
        try:
            return repair_json(raw_output, expected_schema_hint)
        except Exception as second_exc:
            raise BriefingError(
                "Failed to parse LLM JSON output, and JSON repair also failed. "
                f"Initial error: {first_exc}. Repair error: {second_exc}."
            ) from second_exc


def extract_briefing_data(chunks: list[dict[str, Any]]) -> dict[str, Any]:
    """Run NIL preprocessing, merge NIL, select context, and extract final briefing JSON."""
    if not chunks:
        raise BriefingError("No chunks were generated from the input document.")

    nil_results = nil_preprocess(chunks)
    merged_nil = merge_nil_results(nil_results)
    selected_context = _select_context(chunks, merged_nil)
    raw = call_llm(build_final_extraction_prompt(merged_nil, selected_context))
    schema_hint = "Final briefing JSON with paper_title, core_vision, layers, timeline, one_sentence_summary."
    data = parse_or_repair_json(raw, schema_hint)
    _validate_briefing_data(data)
    return data


def _select_context(chunks: list[dict[str, Any]], merged_nil: dict[str, Any]) -> str:
    keywords = [
        "核心愿景",
        "efficiency",
        "intelligence per joule",
        "1000",
        "hardware",
        "algorithm",
        "application",
        "2.2",
        "2-5",
        "6-10",
        "2035",
        "technology trends",
    ]
    scored: list[tuple[int, int, dict[str, Any]]] = []
    for chunk in chunks:
        lower_text = chunk["text"].lower()
        score = sum(3 for kw in keywords if kw.lower() in lower_text)
        score += min(len(lower_text), 4000) // 2000
        scored.append((score, chunk["chunk_id"], chunk))

    if any(score > 0 for score, _, _ in scored):
        selected = [item[2] for item in sorted(scored, key=lambda x: (-x[0], x[1]))[:4]]
    else:
        source_counts: dict[int, int] = {}
        for insight in merged_nil.get("high_confidence_insights", []):
            for chunk_id in insight.get("source_chunk_ids", []):
                source_counts[int(chunk_id)] = source_counts.get(int(chunk_id), 0) + 1
        selected = sorted(
            chunks,
            key=lambda c: (-source_counts.get(c["chunk_id"], 0), c["chunk_id"]),
        )[:4]

    selected = sorted(selected, key=lambda c: c["chunk_id"])
    parts = [
        f"[chunk_id={chunk['chunk_id']}; section={chunk.get('section_title', '')}]\n{chunk['text']}"
        for chunk in selected
    ]
    return "\n\n---\n\n".join(parts)


def _validate_briefing_data(data: dict[str, Any]) -> None:
    required = ["paper_title", "core_vision", "layers", "timeline", "one_sentence_summary"]
    missing = [key for key in required if key not in data]
    if missing:
        raise BriefingError(f"Final briefing JSON is missing required fields: {', '.join(missing)}")

    if not isinstance(data.get("layers"), list) or len(data["layers"]) < 3:
        raise BriefingError("Final briefing JSON must contain three layers.")
    if not isinstance(data.get("timeline"), list) or len(data["timeline"]) < 2:
        raise BriefingError("Final briefing JSON must contain two timeline stages.")


def create_ppt(data: dict[str, Any], output_path: str = "briefing.pptx") -> None:
    """Create a 16:9 three-slide briefing deck."""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        add_title_slide(prs, data)
        add_layers_slide(prs, data)
        add_timeline_slide(prs, data)

        prs.save(output_path)
    except Exception as exc:
        raise BriefingError(f"Failed to generate PPT: {exc}") from exc


def add_title_slide(prs: Presentation, data: dict[str, Any]) -> None:
    """Add title and core vision slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, COLORS["deep_blue"])
    _add_accent_bar(slide)

    title = shorten_text(data.get("paper_title", "论文标题未明确说明"), 95)
    _add_textbox(slide, title, 0.7, 0.55, 7.0, 0.9, 28, COLORS["white"], bold=True)

    summary = shorten_text(data.get("one_sentence_summary", ""), 95)
    _add_textbox(slide, summary, 0.72, 1.45, 6.7, 0.65, 15, RGBColor(0xC7, 0xE5, 0xFF))

    core = data.get("core_vision", {}) or {}
    headline = shorten_text(core.get("headline", "核心愿景未明确说明"), 56)
    highlight = _extract_highlight(headline)

    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(7.65),
        Inches(0.88),
        Inches(4.85),
        Inches(5.38),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["white"]
    card.line.color.rgb = RGBColor(0x99, 0xCC, 0xFF)

    _add_textbox(slide, "CORE VISION", 8.05, 1.18, 3.8, 0.35, 11, COLORS["tech_blue"], bold=True)
    _add_textbox(slide, highlight, 8.0, 1.55, 4.0, 0.78, 30, COLORS["deep_blue"], bold=True)
    _add_textbox(slide, headline, 8.0, 2.35, 4.1, 0.7, 15, COLORS["deep_gray"], bold=True)

    ipj = shorten_text(core.get("intelligence_per_joule", "论文未明确说明"), 100)
    target = shorten_text(core.get("target_explanation", "论文未明确说明"), 130)
    _add_textbox(slide, "Intelligence per Joule", 8.0, 3.22, 4.1, 0.35, 13, COLORS["tech_blue"], bold=True)
    _add_textbox(slide, ipj, 8.0, 3.62, 4.0, 0.85, 13, COLORS["deep_gray"])
    _add_textbox(slide, "为什么重要", 8.0, 4.72, 4.1, 0.35, 13, COLORS["tech_blue"], bold=True)
    _add_textbox(slide, target, 8.0, 5.1, 4.0, 0.75, 12, COLORS["deep_gray"])

    _add_textbox(slide, "AI + Hardware Co-design Briefing", 0.72, 6.75, 5.0, 0.35, 12, RGBColor(0xC7, 0xE5, 0xFF))


def add_layers_slide(prs: Presentation, data: dict[str, Any]) -> None:
    """Add three abstraction-layer cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, COLORS["light_gray"])
    _add_textbox(
        slide,
        "三大抽象层级：Hardware × Algorithm × Application",
        0.55,
        0.42,
        12.0,
        0.55,
        24,
        COLORS["deep_blue"],
        bold=True,
    )

    layers = _normalize_layers(data.get("layers", []))
    x_positions = [0.65, 4.58, 8.51]
    for x, layer in zip(x_positions, layers):
        _add_layer_card(slide, x, 1.35, 3.55, 5.55, layer)


def add_timeline_slide(prs: Presentation, data: dict[str, Any]) -> None:
    """Add horizontal timeline slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, COLORS["white"])
    _add_textbox(slide, "技术演进时间轴：近期到远期", 0.55, 0.42, 12.0, 0.55, 24, COLORS["deep_blue"], bold=True)

    timeline = _normalize_timeline(data.get("timeline", []))
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(1.45),
        Inches(3.23),
        Inches(11.85),
        Inches(3.23),
    )
    line.line.color.rgb = COLORS["tech_blue"]
    line.line.width = Pt(3)
    try:
        line.line.end_arrowhead = True
    except Exception:
        pass

    node_specs = [(3.25, timeline[0]), (9.25, timeline[1])]
    for x, stage in node_specs:
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(2.98), Inches(0.5), Inches(0.5))
        node.fill.solid()
        node.fill.fore_color.rgb = COLORS["tech_blue"]
        node.line.color.rgb = COLORS["white"]
        _add_timeline_card(slide, x - 1.45, 3.75, 3.45, 2.75, stage)

    _add_textbox(slide, "近期", 2.75, 2.45, 1.5, 0.35, 14, COLORS["deep_blue"], bold=True, align=PP_ALIGN.CENTER)
    _add_textbox(slide, "远期", 8.75, 2.45, 1.5, 0.35, 14, COLORS["deep_blue"], bold=True, align=PP_ALIGN.CENTER)


def shorten_text(text: Any, max_len: int) -> str:
    """Shorten text for slide display without mutating source data."""
    value = str(text or "").strip()
    if len(value) <= max_len:
        return value
    return value[: max_len - 1].rstrip() + "…"


def join_bullets(items: Any, max_items: int = 3, max_len_each: int = 40) -> list[str]:
    """Return a short bullet list suitable for a PPT card."""
    if not isinstance(items, list):
        items = [items] if items else []
    bullets = []
    for item in items[:max_items]:
        item_text = shorten_text(item, max_len_each)
        if item_text:
            bullets.append(item_text)
    return bullets


def add_wrapped_textbox(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    font_size: int = 13,
    color: RGBColor = COLORS["deep_gray"],
    bold: bool = False,
):
    """Compatibility helper requested by the project spec."""
    return _add_textbox(slide, text, x, y, w, h, font_size, color, bold=bold)


def _add_background(slide, color: RGBColor) -> None:
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _add_accent_bar(slide) -> None:
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["tech_blue"]
    bar.line.fill.background()


def _add_textbox(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
):
    textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = "Microsoft YaHei"
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color
    return textbox


def _add_bullets(slide, bullets: list[str], x: float, y: float, w: float, h: float, font_size: int = 11):
    textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS["deep_gray"]
        p.space_after = Pt(3)
    return textbox


def _add_layer_card(slide, x: float, y: float, w: float, h: float, layer: dict[str, Any]) -> None:
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["white"]
    card.line.color.rgb = RGBColor(0xD8, 0xE3, 0xF0)

    name = shorten_text(layer.get("name", "论文未明确说明"), 24)
    role = shorten_text(layer.get("role", "论文未明确说明"), 65)
    features = join_bullets(layer.get("key_features", []), 3, 34)
    examples = join_bullets(layer.get("examples", []), 2, 34)

    _add_textbox(slide, name, x + 0.25, y + 0.25, w - 0.5, 0.45, 19, COLORS["deep_blue"], bold=True)
    _add_textbox(slide, role, x + 0.25, y + 0.82, w - 0.5, 0.78, 11, COLORS["mid_gray"])

    _add_textbox(slide, "关键特征", x + 0.25, y + 1.78, w - 0.5, 0.3, 12, COLORS["tech_blue"], bold=True)
    _add_bullets(slide, features, x + 0.35, y + 2.16, w - 0.65, 1.25, 10)

    _add_textbox(slide, "代表性技术或例子", x + 0.25, y + 3.72, w - 0.5, 0.3, 12, COLORS["tech_blue"], bold=True)
    _add_bullets(slide, examples, x + 0.35, y + 4.1, w - 0.65, 0.95, 10)


def _add_timeline_card(slide, x: float, y: float, w: float, h: float, stage: dict[str, Any]) -> None:
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["light_blue"]
    card.line.color.rgb = RGBColor(0xB8, 0xD8, 0xF6)

    title = f"{shorten_text(stage.get('stage', ''), 8)} · {shorten_text(stage.get('range', ''), 16)}"
    trends = join_bullets(stage.get("trends", []), 3, 34)
    meaning = shorten_text(stage.get("meaning", "论文未明确说明"), 70)

    _add_textbox(slide, title, x + 0.25, y + 0.22, w - 0.5, 0.4, 16, COLORS["deep_blue"], bold=True)
    _add_bullets(slide, trends, x + 0.35, y + 0.78, w - 0.7, 0.95, 10)
    _add_textbox(slide, meaning, x + 0.25, y + 1.86, w - 0.5, 0.62, 10, COLORS["deep_gray"])


def _extract_highlight(headline: str) -> str:
    match = re.search(r"(\d+(?:,\d{3})*(?:\.\d+)?\s*[xX倍]?)", headline)
    if match:
        return match.group(1)
    words = re.split(r"[，,。:：；;\s]+", headline)
    candidates = [word for word in words if word]
    return shorten_text(candidates[0] if candidates else headline, 18)


def _normalize_layers(layers: list[dict[str, Any]]) -> list[dict[str, Any]]:
    desired = ["Hardware", "Algorithm", "Application"]
    by_name = {str(layer.get("name", "")).lower(): layer for layer in layers if isinstance(layer, dict)}
    normalized = []
    for name in desired:
        layer = by_name.get(name.lower())
        if not layer:
            layer = {
                "name": name,
                "role": "论文未明确说明",
                "key_features": ["论文未明确说明"],
                "examples": ["论文未明确说明"],
            }
        normalized.append(layer)
    return normalized


def _normalize_timeline(timeline: list[dict[str, Any]]) -> list[dict[str, Any]]:
    defaults = [
        {"stage": "近期", "range": "2-5 年", "trends": ["论文未明确说明"], "meaning": "论文未明确说明"},
        {"stage": "远期", "range": "6-10 年", "trends": ["论文未明确说明"], "meaning": "论文未明确说明"},
    ]
    if not isinstance(timeline, list):
        return defaults
    result = []
    for idx, default in enumerate(defaults):
        item = timeline[idx] if idx < len(timeline) and isinstance(timeline[idx], dict) else {}
        merged = {**default, **item}
        result.append(merged)
    return result


def main() -> None:
    global _LLM_CLIENT, _LLM_MODEL
    parser = argparse.ArgumentParser(description="Generate an AI + hardware co-design briefing deck from a PDF or TXT paper.")
    parser.add_argument("--input", required=True, help="Input paper path, PDF or TXT.")
    parser.add_argument("--output", default="briefing.pptx", help="Output PPTX path. Default: briefing.pptx")
    parser.add_argument("--max-chars", type=int, default=7000, help="Maximum characters per chunk. Default: 7000")
    parser.add_argument("--overlap", type=int, default=300, help="Overlap characters between chunks. Default: 300")
    parser.add_argument("--api-key", help="API key (prefer hidden prompt or an environment variable).")
    parser.add_argument("--base-url", help="OpenAI-compatible API base URL. Defaults to the OpenAI endpoint.")
    parser.add_argument("--model", help="Model name. Default: API_MODEL or gpt-4o-mini.")
    args = parser.parse_args()

    try:
        _LLM_CLIENT, _LLM_MODEL = get_llm_config(args)
        raw_text = load_document(args.input)
        text = clean_text(raw_text)
        if not text:
            raise BriefingError("Input document contains no usable text after cleaning.")
        chunks = chunk_text(text, max_chars=args.max_chars, overlap_chars=args.overlap)
        print(f"Document split into {len(chunks)} chunks.")
        data = extract_briefing_data(chunks)
        create_ppt(data, args.output)
        print(f"{args.output} generated successfully.")
    except BriefingError as exc:
        print(f"Error: {exc}", file=sys.stderr)
        sys.exit(1)
    except KeyboardInterrupt:
        print("Interrupted by user.", file=sys.stderr)
        sys.exit(130)


if __name__ == "__main__":
    main()
